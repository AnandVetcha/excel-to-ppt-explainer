import subprocess
from pptx import Presentation
from pathlib import Path

# Generate the PPTX using sample workbook

def generate_pptx():
    cmd = [
        'python', 'auto_generate_ppt_openpyxl.py',
        '--xlsx', 'sample_sales_mix.xlsx',
        '--sheet', 'Sheet1',
        '--summary_start', 'A12',
        '--key_header', 'Product',
        '--out', 'deck.pptx',
        '--link_mode', 'overlay',
        '--table_font_pt', '12',
        '--round_digits', '2',
        '--skip_cols', '2', '4'
    ]
    subprocess.run(cmd, check=True)

def verify_pptx():
    pptx_path = Path('deck.pptx')
    if not pptx_path.exists():
        raise FileNotFoundError('deck.pptx was not created')

    prs = Presentation(pptx_path)
    expected_slides = 21
    actual_slides = len(prs.slides)
    assert actual_slides == expected_slides, f'Expected {expected_slides} slides, got {actual_slides}'

    first_slide_texts = [
        shape.text for shape in prs.slides[0].shapes if hasattr(shape, 'text')
    ]
    if not any('Summary Table' in t for t in first_slide_texts):
        raise AssertionError('Summary slide does not contain "Summary Table" title')

    print('PPTX verification passed: file exists, has 21 slides, and starts with a Summary Table slide.')

if __name__ == '__main__':
    generate_pptx()
    verify_pptx()
