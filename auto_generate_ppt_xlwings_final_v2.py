
"""
auto_generate_ppt_xlwings_final.py

- Shapes-based hyperlinks are the default (in FRONT of the table).
- True transparency for overlays:
    1) Try python-pptx API: fore_color.transparency = 1.0 (fallback: fill.transparency = 1.0)
    2) Also force the underlying DrawingML: <a:solidFill><a:srgbClr ...><a:alpha val="0"/></a:srgbClr></a:solidFill>
- Overlays are aligned using ACTUAL table widths/heights after text is placed.
- We still add run-level links to the numbers as a backup.
- Table font size is set with --table_font_pt (default 12) and word_wrap=False.
- Numeric values are rounded using --round_digits (default 2).
- (Deprecated alias: --header_font_pt)

Usage:
python auto_generate_ppt_xlwings_final_v2.py  --xlsx sample_sales_mix.xlsx  --sheet Sheet1  --summary_start A12  --raw_table Raw_Data  --key_header Product  --out deck.pptx  --link_mode overlay  --table_font_pt 12  --round_digits 2  --verbose
"""
import argparse
import re
from pathlib import Path

import pandas as pd
import xlwings as xw
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn

# ---------------- Excel helpers ----------------
def get_formula_str(cell):
    for attr in ("Formula2", "Formula"):
        try:
            s = getattr(cell.api, attr)
            if isinstance(s, str) and s.startswith("="):
                return s
        except Exception:
            pass
    return None

def detect_summary_region_from_start(sht, start_addr, max_cols=60, verbose=False):
    start_cell = sht.range(start_addr)
    hdr_row = start_cell.row - 1
    start_col = start_cell.column

    headers, c = [], start_col
    while True:
        val = sht.range((hdr_row, c)).value
        if val in (None, ""):
            break
        headers.append(str(val))
        c += 1
        if c - start_col > max_cols:
            break

    data_rows, r = [], start_cell.row
    while True:
        v = sht.range((r, start_col)).value
        if v in (None, ""):
            break
        data_rows.append(r)
        r += 1

    if verbose:
        print(f"[summary] header row={hdr_row}, headers={headers}, data_rows={data_rows}")
    return hdr_row, headers, data_rows, start_col

def parse_structured_columns(formula, table_name):
    if not formula:
        return []
    cols, target, s, i = [], f"{table_name}[", formula, 0
    while True:
        start = s.find(target, i)
        if start == -1:
            break
        j = start + len(target)
        buf = []
        while j < len(s):
            ch = s[j]
            if ch == ']':
                if j + 1 < len(s) and s[j+1] == ']':
                    buf.append(']'); j += 2; continue
                else:
                    j += 1; break
            else:
                buf.append(ch); j += 1
        name = ''.join(buf).replace("'", "")
        if name not in cols:
            cols.append(name)
        i = j
    return cols

def extract_filter_key(formula, table_name, sht, row_idx, key_col_idx):
    if not formula:
        return (None, None)
    s = formula.replace(" ", "")
    col_letter = xw.utils.col_name(key_col_idx)
    pat = re.compile(
        rf"(?:{re.escape(table_name)}\[([^\]]+?)\]=\$?{col_letter}\$?{row_idx}|\$?{col_letter}\$?{row_idx}={re.escape(table_name)}\[([^\]]+?)\])"
    )
    m = pat.search(s)
    if m:
        col = (m.group(1) or m.group(2) or "").replace("'", "")
        key_value = sht.range((row_idx, key_col_idx)).value
        return (col, key_value)
    return (None, None)

def read_listobject_df(sht, lo_name=None):
    lo = None
    if lo_name:
        try:
            lo = sht.api.ListObjects(lo_name)
        except Exception:
            lo = None
    if lo is None:
        if sht.api.ListObjects.Count >= 1:
            lo = sht.api.ListObjects(1)
        else:
            raise RuntimeError("No Excel Table (ListObject) found on this sheet.")
    addr = lo.Range.Address
    df = sht.range(addr).options(pd.DataFrame, header=1, index=False).value
    return df, lo.Name

def guess_key_col(df_raw, preferred_name):
    import re as _re
    if preferred_name in df_raw.columns:
        return preferred_name
    target_norms = {'subsystem','subsystemname','sub_system', 'sub system'}
    best = None
    for nm in df_raw.columns:
        norm = _re.sub(r'\\W+', ' ', str(nm)).strip().lower()
        if norm in target_norms:
            return nm
        if 'sub' in norm and 'system' in norm:
            best = nm
    return best or df_raw.columns[0]

# ---------------- PPT helpers ----------------

def _force_xml_alpha_zero(shape):
    """
    Force 100% transparency by injecting a:alpha val="0" into the solid fill color node.
    Works with python-pptx's BaseOxmlElement.xpath (no namespaces kwarg).
    """
    el = shape._element
    srgb = el.xpath('.//a:solidFill/a:srgbClr')
    if srgb:
        clr = srgb[0]
    else:
        scheme = el.xpath('.//a:solidFill/a:schemeClr')
        if not scheme:
            return
        clr = scheme[0]
    # remove any existing alpha child
    for a in clr.xpath('./a:alpha'):
        clr.remove(a)
    alpha = OxmlElement('a:alpha')
    alpha.set('val', '0')  # 0 = 0% opacity
    clr.append(alpha)
def add_overlay_link(summary_slide, x, y, w, h, target_slide):
    rect = summary_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    # Make invisible
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(255, 255, 255)
    ok = False
    try:
        rect.fill.fore_color.transparency = 1.0
        ok = True
    except Exception:
        pass
    try:
        rect.fill.transparency = 1.0
        ok = True
    except Exception:
        pass
    # XML-level fallback (always run to be safe)
    _force_xml_alpha_zero(rect)

    rect.line.fill.background()
    rect.line.width = 0
    rect.shadow.inherit = False

    rect.click_action.target_slide = target_slide
    return rect

# ---------------- Formatting helper ----------------
def format_number(val, round_digits: int) -> str:
    if val is None:
        return ""
    if isinstance(val, (int, float)) and not isinstance(val, bool):
        fmt = f"{{:.{round_digits}f}}"
        return fmt.format(val)
    return str(val)

# ---------------- Builder ----------------
def build_ppt_xlwings(xlsx_path: Path, out_path: Path, sheet_name: str, summary_start: str, raw_table_name: str=None, verbose: bool=False, link_mode: str="overlay", table_font_pt: int=12, key_header: str=None, round_digits: int=2):
    app = xw.App(visible=False, add_book=False)
    try:
        wb = xw.Book(xlsx_path)
        sht = wb.sheets[sheet_name] if sheet_name else wb.sheets.active
        wb.app.api.CalculateFull()

        df_raw, actual_table_name = read_listobject_df(sht, raw_table_name)
        raw_table_name = actual_table_name

        hdr_row, headers, data_rows, start_col_idx = detect_summary_region_from_start(sht, summary_start, verbose=verbose)
        if not headers or not data_rows:
            raise RuntimeError("Could not detect headers or data rows; check --summary_start.")

        if key_header is None:
            key_header = str(headers[0])

        # collect formulas/values
        last_row_in_block = data_rows[-1]
        summary = []
        for r in data_rows:
            key_value = sht.range((r, start_col_idx)).value
            items = {"row": r, "key": key_value, "cells": {}}
            for c_off, h in enumerate(headers[1:], start=1):
                c_idx = start_col_idx + c_off
                rng = sht.range((r, c_idx))
                f = get_formula_str(rng)
                if not f:
                    rr = r - 1
                    while rr >= hdr_row + 1 and not f:
                        f = get_formula_str(sht.range((rr, c_idx))); rr -= 1
                    if not f:
                        rr = r + 1
                        while rr <= last_row_in_block and not f:
                            f = get_formula_str(sht.range((rr, c_idx))); rr += 1
                items["cells"][h] = {"address": rng.get_address(), "formula": f, "value": rng.value}
                if verbose:
                    print(f"[cell] r={r}, c_idx={c_idx}, header={h}, formula_found={bool(f)}")
            summary.append(items)

        prs = Presentation()
        # Title Only layout
        summary_slide = prs.slides.add_slide(prs.slide_layouts[5])
        summary_slide.shapes.title.text = "Summary Table"

        # table scaffold
        sum_cols = len(headers)
        sum_rows = len(summary) + 1
        left, top, width = Inches(0.5), Inches(1.5), Inches(9.5)
        # Derive a row height that scales with the font size.
        # 0.4" works well for an 18pt font, so convert 0.4" to points and
        # use that to compute a proportional height for the requested font size.
        points_per_inch = 72
        row_height_pt = (0.4 * points_per_inch / 18) * table_font_pt
        base_row_height = Pt(row_height_pt)

        table_shape = summary_slide.shapes.add_table(
            sum_rows, sum_cols, left, top, width, base_row_height * sum_rows
        )
        table = table_shape.table

        total_w = int(width)
        base_w = total_w // sum_cols
        remainder = total_w - base_w * sum_cols
        col_widths = [base_w] * sum_cols
        col_widths[-1] = base_w + remainder
        for j in range(sum_cols):
            table.columns[j].width = col_widths[j]
        for i in range(sum_rows):
            table.rows[i].height = int(base_row_height)

        # header text (no wrap)
        for j, h in enumerate(headers):
            tf = table.cell(0, j).text_frame
            tf.clear()
            tf.word_wrap = False
            run = tf.paragraphs[0].add_run()
            run.text = str(h)
            run.font.bold = True
            run.font.size = Pt(table_font_pt)

        # build detail slides first
        detail_slide_map = {}
        for i, row in enumerate(summary, start=1):
            key = row["key"]
            for j, metric in enumerate(headers[1:], start=1):
                info = row["cells"][metric]
                formula = info["formula"]
                cols_used = [key_header] + parse_structured_columns(formula, raw_table_name)
                cols_used = [c for c in cols_used if c in df_raw.columns]
                if not cols_used:
                    cols_used = [key_header] if key_header in df_raw.columns else list(df_raw.columns)
                colname, key_from_formula = extract_filter_key(formula, raw_table_name, sht, row["row"], key_col_idx=start_col_idx)
                if colname is None:
                    colname = guess_key_col(df_raw, key_header)
                key_val = key_from_formula if key_from_formula is not None else key
                try:
                    df_filtered = df_raw[df_raw[colname] == key_val]
                except Exception:
                    key_col = guess_key_col(df_raw, key_header)
                    df_filtered = df_raw[df_raw[key_col] == key_val]
                df_snippet = df_filtered[cols_used].copy()

                slide = prs.slides.add_slide(prs.slide_layouts[5])
                slide.shapes.title.text = f"{key} – {metric}"
                # Home button to return to summary
                btn = slide.shapes.add_shape(
                    MSO_SHAPE.ACTION_BUTTON_HOME,
                    Inches(9.0),
                    Inches(0.2),
                    Inches(0.5),
                    Inches(0.5),
                )
                btn.click_action.target_slide = summary_slide
                btn.text_frame.text = ""
                # Formula box
                tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9.2), Inches(1.2))
                tf = tx.text_frame; tf.clear()
                tf.word_wrap = True
                p1 = tf.paragraphs[0]; p1.text = "Formula:"; p1.font.bold = True
                p2 = tf.add_paragraph(); p2.text = formula if formula else "(no formula found)"; p2.level = 1; p2.font.size = Pt(14) 
                p3 = tf.add_paragraph(); p3.text = f"Evaluated value: {format_number(info['value'], round_digits)}"; p3.level = 1;p3.font.size = Pt(14)
                # Snippet
                rows, cols = df_snippet.shape
                s_table = slide.shapes.add_table(rows+1, cols, Inches(0.5), Inches(2.6), Inches(9.2), Inches(0.6 + 0.3*max(rows,1))).table
                for jj, hh in enumerate(df_snippet.columns):
                    tfh = s_table.cell(0, jj).text_frame; tfh.clear()
                    r0 = tfh.paragraphs[0].add_run(); r0.text = str(hh); r0.font.bold = True; r0.font.size = Pt(table_font_pt)
                for ii in range(rows):
                    for jj in range(cols):
                        val = df_snippet.iloc[ii, jj]
                        cell = s_table.cell(ii+1, jj)
                        tfcell = cell.text_frame; tfcell.clear()
                        run = tfcell.paragraphs[0].add_run()
                        run.text = format_number(val, round_digits)
                        run.font.size = Pt(table_font_pt)
                detail_slide_map[(i, metric)] = slide

        # write summary values
        for i, row in enumerate(summary, start=1):
            tf0 = table.cell(i, 0).text_frame; tf0.clear()
            run0 = tf0.paragraphs[0].add_run()
            run0.text = format_number(row['key'], round_digits)
            run0.font.size = Pt(table_font_pt)
            for j, metric in enumerate(headers[1:], start=1):
                tf = table.cell(i, j).text_frame; tf.clear()
                run = tf.paragraphs[0].add_run()
                val = row["cells"][metric]["value"]
                text = format_number(val, round_digits)
                run.text = text
                run.font.size = Pt(table_font_pt)
                target = detail_slide_map.get((i, metric))
                if target and text != "":
                    run.hyperlink.target_slide = target  # backup link

        # recompute actual grid
        table_left = int(table_shape.left)
        table_top = int(table_shape.top)

        col_lefts = [table_left]
        for j in range(1, sum_cols):
            col_lefts.append(col_lefts[-1] + int(table.columns[j - 1].width))

        row_tops = [table_top]
        for i in range(sum_rows):
            row_tops.append(row_tops[-1] + int(table.rows[i].height))

        # overlays in FRONT
        if link_mode == "overlay":
            for i in range(1, sum_rows):
                for j, metric in enumerate(headers[1:], start=1):
                    target = detail_slide_map.get((i, metric))
                    if not target:
                        continue
                    cell = table.cell(i, j)
                    txt = cell.text_frame.paragraphs[0].runs[0].text if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs else ""
                    if txt == "":
                        continue
                    x = col_lefts[j]
                    y = row_tops[i]
                    w = int(table.columns[j].width)
                    h = row_tops[i + 1] - row_tops[i]
                    add_overlay_link(summary_slide, x, y, w, h, target)

        prs.save(out_path)
        return out_path
    finally:
        try:
            wb.close()
        except Exception:
            pass
        app.quit()

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--xlsx", required=True, help="Path to Excel file")
    ap.add_argument("--sheet", default="Sheet1", help="Worksheet name")
    ap.add_argument("--summary_start", required=True, help="Top-left data cell of summary (e.g., A12)")
    ap.add_argument("--out", default="deck.pptx", help="Output PPTX")
    ap.add_argument("--raw_table", default=None, help="Excel Table (ListObject) name (optional)")
    ap.add_argument("--key_header", default=None, help="Column to display in detail tables (e.g., 'Product')")
    ap.add_argument("--link_mode", choices=["text","overlay"], default="overlay", help="How to create links on summary cells")
    ap.add_argument("--table_font_pt", type=int, default=None, help="Font size for table text")
    ap.add_argument("--header_font_pt", type=int, default=None, help=argparse.SUPPRESS)
    ap.add_argument("--round_digits", type=int, default=2, help="Decimal places for numeric values")
    ap.add_argument("--verbose", action="store_true", help="Debug prints")
    args = ap.parse_args()
    font_pt = args.table_font_pt if args.table_font_pt is not None else args.header_font_pt
    if font_pt is None:
        font_pt = 12

    out = build_ppt_xlwings(
        xlsx_path=Path(args.xlsx),
        out_path=Path(args.out),
        sheet_name=args.sheet,
        summary_start=args.summary_start,
        raw_table_name=args.raw_table,
        verbose=args.verbose,
        link_mode=args.link_mode,
        table_font_pt=font_pt,
        key_header=args.key_header,
        round_digits=args.round_digits,
    )
    print(f"PPT created: {out}")

if __name__ == "__main__":
    main()
