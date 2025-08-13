
"""
auto_generate_ppt_openpyxl.py

- Shapes-based hyperlinks are the default (in FRONT of the table).
- True transparency for overlays:
    1) Try python-pptx API: fore_color.transparency = 1.0 (fallback: fill.transparency = 1.0)
    2) Also force the underlying DrawingML: <a:solidFill><a:srgbClr ...><a:alpha val="0"/></a:srgbClr></a:solidFill>
- Overlays are aligned using ACTUAL table widths/heights after text is placed.
- We still add run-level links to the numbers as a backup.
- Table font size is set with --table_font_pt (default 12) and word_wrap=False.
- Numeric values are rounded using --round_digits (default 2).
- (Deprecated alias: --header_font_pt)

Usage:
python auto_generate_ppt_openpyxl.py  --xlsx sample_sales_mix.xlsx  --sheet Sheet1  --summary_start A12  --key_header Product  --out deck.pptx  --link_mode overlay  --table_font_pt 12  --round_digits 2  --skip_cols 2 4  --verbose
"""
import argparse
import re
from pathlib import Path

import pandas as pd
from openpyxl import load_workbook
from openpyxl.utils.cell import coordinate_to_tuple, get_column_letter, range_boundaries
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

# ---------------- Excel helpers ----------------
def get_formula_str(cell):
    val = cell.value
    # openpyxl marks formulas with data_type "f", but depending on version
    # the value may or may not include the leading '='. Normalize here so
    # callers consistently receive formulas prefixed with '='.
    if getattr(cell, "data_type", None) == "f":
        text = "" if val is None else str(val)
        return text if text.startswith("=") else f"={text}"
    if isinstance(val, str) and val.startswith("="):
        return val
    return None

def detect_summary_region_from_start(ws, start_addr, max_cols=60, verbose=False):
    start_row, start_col = coordinate_to_tuple(start_addr)
    hdr_row = start_row - 1

    headers, c = [], start_col
    while True:
        val = ws.cell(row=hdr_row, column=c).value
        if val in (None, ""):
            break
        headers.append(str(val))
        c += 1
        if c - start_col > max_cols:
            break

    data_rows, r = [], start_row
    while True:
        v = ws.cell(row=r, column=start_col).value
        if v in (None, ""):
            break
        data_rows.append(r)
        r += 1

    if verbose:
        print(f"[summary] header row={hdr_row}, headers={headers}, data_rows={data_rows}")
    return hdr_row, headers, data_rows, start_col

def parse_structured_columns(formula, table_name):
    if not formula:
        return []
    cols, target, s, i = [], f"{table_name}[", formula, 0
    while True:
        start = s.find(target, i)
        if start == -1:
            break
        j = start + len(target)
        buf = []
        while j < len(s):
            ch = s[j]
            if ch == ']':
                if j + 1 < len(s) and s[j+1] == ']':
                    buf.append(']'); j += 2; continue
                else:
                    j += 1; break
            else:
                buf.append(ch); j += 1
        name = ''.join(buf).replace("'", "")
        if name not in cols:
            cols.append(name)
        i = j
    return cols

def extract_table_names(formula):
    if not formula:
        return []
    s = formula.replace("'", "")
    names = re.findall(r"([A-Za-z0-9_]+)\[", s)
    seen = []
    for n in names:
        if n not in seen:
            seen.append(n)
    return seen

def extract_filter_key(formula, table_name, sht, row_idx, key_col_idx):
    if not formula:
        return (None, None)
    s = formula.replace(" ", "")
    col_letter = get_column_letter(key_col_idx)
    cell_pat = rf"\$?{col_letter}\$?{row_idx}"
    # Equality comparison (e.g., Table1[Product]=$A12)
    pat_eq = re.compile(
        rf"(?:{re.escape(table_name)}\[([^\]]+?)\]={cell_pat}|{cell_pat}={re.escape(table_name)}\[([^\]]+?)\])"
    )
    m = pat_eq.search(s)
    if m:
        col = (m.group(1) or m.group(2) or "").replace("'", "")
        key_value = sht.cell(row=row_idx, column=key_col_idx).value
        return (col, key_value)
    # Function-style criteria (e.g., SUMIFS(Table1[Amount],Table1[Product],$A12))
    pat_func = re.compile(rf"{re.escape(table_name)}\[([^\]]+?)\],{cell_pat}")
    m = pat_func.search(s)
    if m:
        col = m.group(1).replace("'", "")
        key_value = sht.cell(row=row_idx, column=key_col_idx).value
        return (col, key_value)
    pat_func_rev = re.compile(rf"{cell_pat},{re.escape(table_name)}\[([^\]]+?)\]")
    m = pat_func_rev.search(s)
    if m:
        col = m.group(1).replace("'", "")
        key_value = sht.cell(row=row_idx, column=key_col_idx).value
        return (col, key_value)
    return (None, None)
def read_all_tables(wb_formula, wb_values, verbose: bool = False):
    """Return DataFrames and column formulas for all Excel Tables."""
    tables = {}
    table_formulas = {}
    for ws_formula in wb_formula.worksheets:
        ws_values = wb_values[ws_formula.title]
        if not ws_formula.tables:
            continue
        if hasattr(ws_formula.tables, "items"):
            table_iter = ws_formula.tables.items()
        else:
            table_iter = [(t.name, t) for t in ws_formula.tables]
        for name, tbl in table_iter:
            ref = getattr(tbl, "ref", tbl)
            min_col, min_row, max_col, max_row = range_boundaries(ref)
            headers = [ws_values.cell(row=min_row, column=c).value for c in range(min_col, max_col + 1)]
            data = []
            formula_rows = []
            for r in range(min_row + 1, max_row + 1):
                row_vals = []
                row_fml = []
                for c in range(min_col, max_col + 1):
                    row_vals.append(ws_values.cell(row=r, column=c).value)
                    row_fml.append(get_formula_str(ws_formula.cell(row=r, column=c)))
                data.append(row_vals)
                formula_rows.append(row_fml)
            df = pd.DataFrame(data, columns=headers)
            df_fml = pd.DataFrame(formula_rows, columns=headers)
            tables[name] = df
            table_formulas[name] = df_fml
            if verbose:
                for j, h in enumerate(headers):
                    formula = df_fml.iloc[0, j]
                    parsed = parse_structured_columns(formula, name) if formula else []
                    print(f"[table] {name} column={h} formula={formula} parsed={parsed}")
    if not tables:
        raise RuntimeError("No Excel Table (ListObject) found in this workbook.")
    return tables, table_formulas

def guess_key_col(df_raw, preferred_name):
    import re as _re
    if preferred_name in df_raw.columns:
        return preferred_name
    target_norms = {'subsystem','subsystemname','sub_system', 'sub system'}
    best = None
    for nm in df_raw.columns:
        norm = _re.sub(r'\\W+', ' ', str(nm)).strip().lower()
        if norm in target_norms:
            return nm
        if 'sub' in norm and 'system' in norm:
            best = nm
    return best or df_raw.columns[0]

# ---------------- PPT helpers ----------------

def _force_xml_alpha_zero(shape):
    """
    Force 100% transparency by injecting a:alpha val="0" into the solid fill color node.
    Works with python-pptx's BaseOxmlElement.xpath (no namespaces kwarg).
    """
    el = shape._element
    srgb = el.xpath('.//a:solidFill/a:srgbClr')
    if srgb:
        clr = srgb[0]
    else:
        scheme = el.xpath('.//a:solidFill/a:schemeClr')
        if not scheme:
            return
        clr = scheme[0]
    # remove any existing alpha child
    for a in clr.xpath('./a:alpha'):
        clr.remove(a)
    alpha = OxmlElement('a:alpha')
    alpha.set('val', '0')  # 0 = 0% opacity
    clr.append(alpha)
def add_overlay_link(summary_slide, x, y, w, h, target_slide):
    rect = summary_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    # Make invisible
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(255, 255, 255)
    ok = False
    try:
        rect.fill.fore_color.transparency = 1.0
        ok = True
    except Exception:
        pass
    try:
        rect.fill.transparency = 1.0
        ok = True
    except Exception:
        pass
    # XML-level fallback (always run to be safe)
    _force_xml_alpha_zero(rect)

    rect.line.fill.background()
    rect.line.width = 0
    rect.shadow.inherit = False

    rect.click_action.target_slide = target_slide
    return rect


def link_run_to_slide(run, dest_slide, tooltip_text: str = ""):
    """Attach an internal hyperlink to a run that jumps to ``dest_slide``.

    Parameters
    ----------
    run : pptx.text.text._Run
        Run to link.
    dest_slide : pptx.slide.Slide
        Slide to jump to when the run is clicked.
    tooltip_text : str, optional
        Tooltip to display on hover.
    """
    rId = run.part.relate_to(dest_slide.part, RT.SLIDE)
    rPr = run._r.get_or_add_rPr()
    # remove any existing run-level hyperlink
    for child in list(rPr):
        if child.tag.endswith("hlinkClick"):
            rPr.remove(child)
    h = OxmlElement("a:hlinkClick")
    h.set(qn("r:id"), rId)
    h.set("action", "ppaction://hlinksldjump")
    if tooltip_text:
        h.set("tooltip", tooltip_text)
    rPr.append(h)

# ---------------- Formatting helper ----------------
def format_number(val, round_digits: int) -> str:
    if val is None:
        return ""
    if isinstance(val, (int, float)) and not isinstance(val, bool):
        fmt = f"{{:.{round_digits}f}}"
        return fmt.format(val)
    return str(val)

# ---------------- Builder ----------------
def build_ppt_openpyxl(
    xlsx_path: Path,
    out_path: Path,
    sheet_name: str,
    summary_start: str,
    raw_table_name: str = None,
    verbose: bool = False,
    link_mode: str = "overlay",
    table_font_pt: int = 12,
    key_header: str = None,
    round_digits: int = 2,
    pptx_in_path: Path = None,
    skip_col_idxs: list[int] | None = None,
):
    wb_formula = load_workbook(xlsx_path, data_only=False)
    wb_values = load_workbook(xlsx_path, data_only=True)
    try:
        ws_formula = wb_formula[sheet_name] if sheet_name else wb_formula.active
        ws_values = wb_values[sheet_name] if sheet_name else wb_values.active

        table_dfs, table_formulas = read_all_tables(wb_formula, wb_values, verbose=verbose)
        if raw_table_name and raw_table_name in table_dfs:
            default_table_name = raw_table_name
        else:
            default_table_name = next(iter(table_dfs))

        hdr_row, headers, data_rows, start_col_idx = detect_summary_region_from_start(ws_values, summary_start, verbose=verbose)
        if not headers or not data_rows:
            raise RuntimeError("Could not detect headers or data rows; check --summary_start.")

        if key_header is None:
            key_header = str(headers[0])

        # collect formulas/values
        last_row_in_block = data_rows[-1]
        summary = []
        for r in data_rows:
            key_value = ws_values.cell(row=r, column=start_col_idx).value
            items = {"row": r, "key": key_value, "cells": {}}
            for c_off, h in enumerate(headers[1:], start=1):
                c_idx = start_col_idx + c_off
                cell_formula = ws_formula.cell(row=r, column=c_idx)
                f = get_formula_str(cell_formula)
                if not f:
                    rr = r - 1
                    while rr >= hdr_row + 1 and not f:
                        f = get_formula_str(ws_formula.cell(row=rr, column=c_idx)); rr -= 1
                    if not f:
                        rr = r + 1
                        while rr <= last_row_in_block and not f:
                            f = get_formula_str(ws_formula.cell(row=rr, column=c_idx)); rr += 1
                tbls = extract_table_names(f)
                val = ws_values.cell(row=r, column=c_idx).value
                items["cells"][h] = {
                    "address": cell_formula.coordinate,
                    "formula": f,
                    "value": val,
                    "table": tbls[0] if tbls else None,
                }
                if verbose:
                    print(f"[cell] r={r}, c_idx={c_idx}, header={h}, formula_found={bool(f)}")
            summary.append(items)

        prs = Presentation(pptx_in_path) if pptx_in_path else Presentation()
        skip_set = set(skip_col_idxs or [])
        # Title Only layout
        summary_slide = prs.slides.add_slide(prs.slide_layouts[5])
        summary_slide.shapes.title.text = "Summary Table"

        # table scaffold aligned with title margins
        sum_cols = len(headers)
        sum_rows = len(summary) + 1
        title_shape = summary_slide.shapes.title
        left = title_shape.left
        right_margin = prs.slide_width - (title_shape.left + title_shape.width)
        top = title_shape.top + title_shape.height + Inches(0.2)
        width = prs.slide_width - left - right_margin
        if link_mode == "text":
            base_row_height = Inches(0.4 * table_font_pt / 18)
        else:
            base_row_height = Inches(0.4)
        table_shape = summary_slide.shapes.add_table(sum_rows, sum_cols, left, top, width, base_row_height * sum_rows)
        table = table_shape.table

        total_w = int(width)
        base_w = total_w // sum_cols
        remainder = total_w - base_w * sum_cols
        col_widths = [base_w] * sum_cols
        col_widths[-1] = base_w + remainder
        for j in range(sum_cols):
            table.columns[j].width = col_widths[j]
        for i in range(sum_rows):
            table.rows[i].height = int(base_row_height)

        # header text (no wrap)
        for j, h in enumerate(headers):
            tf = table.cell(0, j).text_frame
            tf.clear()
            tf.word_wrap = False
            run = tf.paragraphs[0].add_run()
            run.text = str(h)
            run.font.bold = True
            run.font.size = Pt(table_font_pt)

        # build detail slides first
        detail_slide_map = {}
        for i, row in enumerate(summary, start=1):
            key = row["key"]
            for j, metric in enumerate(headers[1:], start=1):
                if j in skip_set:
                    continue
                info = row["cells"][metric]
                formula = info["formula"]
                tbl_name = info.get("table") or default_table_name
                df_raw = table_dfs.get(tbl_name, table_dfs[default_table_name])
                cols_used = [key_header] + parse_structured_columns(formula, tbl_name)
                cols_used = list(dict.fromkeys(cols_used))
                cols_used = [c for c in cols_used if c in df_raw.columns]
                if not cols_used:
                    cols_used = [key_header] if key_header in df_raw.columns else list(df_raw.columns)
                colname, key_from_formula = extract_filter_key(formula, tbl_name, ws_values, row["row"], key_col_idx=start_col_idx)
                if colname is None:
                    colname = guess_key_col(df_raw, key_header)
                key_val = key_from_formula if key_from_formula is not None else key
                try:
                    df_filtered = df_raw[df_raw[colname] == key_val]
                except Exception:
                    key_col = guess_key_col(df_raw, key_header)
                    df_filtered = df_raw[df_raw[key_col] == key_val]
                df_snippet = df_filtered[cols_used].copy()

                slide = prs.slides.add_slide(prs.slide_layouts[5])
                slide.shapes.title.text = f"{key} – {metric}"
                title_shape = slide.shapes.title
                right_margin = prs.slide_width - (title_shape.left + title_shape.width)
                content_left = title_shape.left
                content_width = prs.slide_width - content_left - right_margin
                # Home button to return to summary
                btn_left = prs.slide_width - right_margin - Inches(0.5)
                btn = slide.shapes.add_shape(
                    MSO_SHAPE.ACTION_BUTTON_HOME,
                    btn_left,
                    Inches(0.2),
                    Inches(0.5),
                    Inches(0.5),
                )
                btn.click_action.target_slide = summary_slide
                btn.text_frame.text = ""
                # Formula box
                formula_height = Inches(1.2)
                formula_top = title_shape.top + title_shape.height + Inches(0.2)
                tx = slide.shapes.add_textbox(content_left, formula_top, content_width, formula_height)
                tf = tx.text_frame; tf.clear()
                tf.word_wrap = True
                p1 = tf.paragraphs[0]; p1.text = "Formula:"; p1.font.bold = True
                p2 = tf.add_paragraph(); p2.text = formula if formula else "(no formula found)"; p2.level = 1; p2.font.size = Pt(14)
                p3 = tf.add_paragraph(); p3.text = f"Evaluated value: {format_number(info['value'], round_digits)}"; p3.level = 1;p3.font.size = Pt(14)
                # Snippet
                rows, cols = df_snippet.shape
                if link_mode == "text":
                    snip_row_height = Inches(0.4 * table_font_pt / 18)
                    snip_height = snip_row_height * (rows + 1)
                else:
                    snip_height = Inches(0.6 + 0.3*max(rows,1))
                    snip_row_height = None
                snip_top = formula_top + formula_height + Inches(0.2)
                s_table_shape = slide.shapes.add_table(rows+1, cols, content_left, snip_top, content_width, snip_height)
                s_table = s_table_shape.table
                if snip_row_height is not None:
                    for rr in range(rows+1):
                        s_table.rows[rr].height = int(snip_row_height)
                for jj, hh in enumerate(df_snippet.columns):
                    tfh = s_table.cell(0, jj).text_frame; tfh.clear()
                    r0 = tfh.paragraphs[0].add_run(); r0.text = str(hh); r0.font.bold = True; r0.font.size = Pt(table_font_pt)
                for ii in range(rows):
                    for jj in range(cols):
                        val = df_snippet.iloc[ii, jj]
                        cell = s_table.cell(ii+1, jj)
                        tfcell = cell.text_frame; tfcell.clear()
                        run = tfcell.paragraphs[0].add_run()
                        run.text = format_number(val, round_digits)
                        run.font.size = Pt(table_font_pt)
                detail_slide_map[(i, metric)] = slide

        # write summary values
        for i, row in enumerate(summary, start=1):
            tf0 = table.cell(i, 0).text_frame; tf0.clear()
            run0 = tf0.paragraphs[0].add_run()
            run0.text = format_number(row['key'], round_digits)
            run0.font.size = Pt(table_font_pt)
            for j, metric in enumerate(headers[1:], start=1):
                tf = table.cell(i, j).text_frame; tf.clear()
                run = tf.paragraphs[0].add_run()
                val = row["cells"][metric]["value"]
                text = format_number(val, round_digits)
                run.text = text
                run.font.size = Pt(table_font_pt)
                if j in skip_set:
                    continue
                target = detail_slide_map.get((i, metric))
                if target and text != "":
                    tooltip = target.shapes.title.text if target.shapes.title else ""
                    link_run_to_slide(run, target, tooltip_text=tooltip)

        # recompute actual grid
        col_lefts = [int(left)]
        for j in range(1, sum_cols):
            col_lefts.append(col_lefts[-1] + int(table.columns[j-1].width))
        row_tops = [int(top)]
        for i in range(1, sum_rows):
            row_tops.append(row_tops[-1] + int(table.rows[i-1].height))

        # overlays in FRONT
        if link_mode == "overlay":
            for i in range(1, sum_rows):
                for j, metric in enumerate(headers[1:], start=1):
                    if j in skip_set:
                        continue
                    target = detail_slide_map.get((i, metric))
                    if not target:
                        continue
                    cell = table.cell(i, j)
                    txt = cell.text_frame.paragraphs[0].runs[0].text if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs else ""
                    if txt == "":
                        continue
                    x = col_lefts[j]
                    y = row_tops[i]
                    w = int(table.columns[j].width)
                    h = int(table.rows[i].height)
                    add_overlay_link(summary_slide, x, y, w, h, target)

        prs.save(out_path)
        return out_path
    finally:
        wb_formula.close()
        wb_values.close()

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--xlsx", required=True, help="Path to Excel file")
    ap.add_argument("--sheet", default="Sheet1", help="Worksheet containing the summary table")
    ap.add_argument("--summary_start", required=True, help="Top-left data cell of summary (e.g., A12)")
    ap.add_argument("--pptx_in", default=None, help="Existing PPTX to append slides to")
    ap.add_argument("--out", default=None, help="Output PPTX (defaults to --pptx_in or 'deck.pptx')")
    ap.add_argument("--raw_table", default=None, help="Default Excel Table (ListObject) name (optional; auto-detected if omitted)")
    ap.add_argument("--key_header", default=None, help="Column to display in detail tables (e.g., 'Product')")
    ap.add_argument("--link_mode", choices=["text","overlay"], default="text", help="How to create links on summary cells")
    ap.add_argument("--table_font_pt", type=int, default=None, help="Font size for table text")
    ap.add_argument("--header_font_pt", type=int, default=None, help=argparse.SUPPRESS)
    ap.add_argument("--round_digits", type=int, default=2, help="Decimal places for numeric values")
    ap.add_argument("--verbose", action="store_true", help="Debug prints")
    ap.add_argument(
        "--skip_cols",
        type=int,
        nargs="*",
        default=[],
        help="1-based indices of data columns (excluding the key column) to skip linking",
    )
    args = ap.parse_args()
    font_pt = args.table_font_pt if args.table_font_pt is not None else args.header_font_pt
    if font_pt is None:
        font_pt = 12

    out_path = Path(args.out) if args.out else Path(args.pptx_in) if args.pptx_in else Path("deck.pptx")

    out = build_ppt_openpyxl(
        xlsx_path=Path(args.xlsx),
        out_path=out_path,
        sheet_name=args.sheet,
        summary_start=args.summary_start,
        raw_table_name=args.raw_table,
        verbose=args.verbose,
        link_mode=args.link_mode,
        table_font_pt=font_pt,
        key_header=args.key_header,
        round_digits=args.round_digits,
        pptx_in_path=Path(args.pptx_in) if args.pptx_in else None,
        skip_col_idxs=args.skip_cols,
    )
    print(f"PPT created: {out}")

if __name__ == "__main__":
    main()
